#!/usr/bin/env python3
"""Inspect PowerPoint presentation structure, metadata, and statistics."""

import json
import sys
import os


def ensure_dependencies():
    try:
        import pptx
    except ImportError:
        os.system("pip install python-pptx")


def inspect_presentation(prs_path):
    """Get comprehensive structural information about a .pptx file."""
    ensure_dependencies()
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(prs_path)

    # Metadata
    props = prs.core_properties
    metadata = {
        "title": props.title or "",
        "author": props.author or "",
        "subject": props.subject or "",
        "keywords": props.keywords or "",
        "created": str(props.created) if props.created else None,
        "modified": str(props.modified) if props.modified else None,
        "last_modified_by": props.last_modified_by or "",
        "revision": props.revision,
        "category": props.category or "",
        "comments": props.comments or "",
    }

    # Dimensions
    dimensions = {
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "aspect_ratio": f"{round(prs.slide_width / prs.slide_height, 2)}:1",
    }

    # Common aspect ratios
    ratio = prs.slide_width / prs.slide_height
    if abs(ratio - 16/9) < 0.05:
        dimensions["aspect_ratio_name"] = "16:9 (Widescreen)"
    elif abs(ratio - 4/3) < 0.05:
        dimensions["aspect_ratio_name"] = "4:3 (Standard)"
    elif abs(ratio - 16/10) < 0.05:
        dimensions["aspect_ratio_name"] = "16:10"

    # Per-slide details
    slides_info = []
    total_word_count = 0
    total_char_count = 0
    total_images = 0
    total_tables = 0
    total_shapes = 0

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i,
            "slide_number": i + 1,
            "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "title": slide.shapes.title.text if slide.shapes.title else None,
            "shape_count": len(slide.shapes),
            "has_notes": False,
            "has_table": False,
            "has_image": False,
            "word_count": 0,
        }

        for shape in slide.shapes:
            total_shapes += 1

            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    words = len(p.text.split())
                    chars = len(p.text)
                    slide_info["word_count"] += words
                    total_word_count += words
                    total_char_count += chars

            if shape.has_table:
                slide_info["has_table"] = True
                total_tables += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        words = len(cell.text.split())
                        total_word_count += words

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_info["has_image"] = True
                total_images += 1

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_info["has_notes"] = True
                slide_info["notes_preview"] = notes_text[:100] + ("..." if len(notes_text) > 100 else "")

        slides_info.append(slide_info)

    # Statistics
    statistics = {
        "slide_count": len(prs.slides),
        "total_word_count": total_word_count,
        "total_character_count": total_char_count,
        "total_shapes": total_shapes,
        "total_images": total_images,
        "total_tables": total_tables,
        "slides_with_notes": sum(1 for s in slides_info if s["has_notes"]),
        "slides_with_tables": sum(1 for s in slides_info if s["has_table"]),
        "slides_with_images": sum(1 for s in slides_info if s["has_image"]),
    }

    # Slide layouts available
    layouts = []
    for layout in prs.slide_layouts:
        layouts.append(layout.name)

    # Slide masters
    masters = []
    for master in prs.slide_masters:
        masters.append({
            "name": master.name if hasattr(master, 'name') else "Default",
            "layout_count": len(master.slide_layouts),
        })

    # File info
    file_size = os.path.getsize(prs_path)
    file_info = {
        "path": prs_path,
        "size_bytes": file_size,
        "size_readable": _format_size(file_size),
    }

    return {
        "file": file_info,
        "metadata": metadata,
        "dimensions": dimensions,
        "statistics": statistics,
        "slides": slides_info,
        "available_layouts": layouts,
        "slide_masters": masters,
    }


def _format_size(size_bytes):
    """Format byte size to human-readable string."""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} TB"


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: inspect_pptx.py <file.pptx>"}))
        sys.exit(1)

    prs_path = sys.argv[1]

    if not os.path.exists(prs_path):
        print(json.dumps({"error": f"File not found: {prs_path}"}))
        sys.exit(1)

    try:
        result = inspect_presentation(prs_path)
        print(json.dumps(result, indent=2, default=str))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)


if __name__ == '__main__':
    main()
