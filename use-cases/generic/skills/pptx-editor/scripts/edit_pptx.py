#!/usr/bin/env python3
"""Edit existing PowerPoint presentations (.pptx) with various operations."""

import json
import sys
import os
import copy


def ensure_dependencies():
    try:
        import pptx
    except ImportError:
        os.system("pip install python-pptx")


def find_replace(prs, find_text, replace_text):
    """Find and replace text across all slides."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if find_text in run.text:
                            run.text = run.text.replace(find_text, replace_text)
                            count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if find_text in run.text:
                                    run.text = run.text.replace(find_text, replace_text)
                                    count += 1
        # Notes
        if slide.has_notes_slide:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                for run in paragraph.runs:
                    if find_text in run.text:
                        run.text = run.text.replace(find_text, replace_text)
                        count += 1
    return count


def delete_slide(prs, slide_index):
    """Delete a slide by index."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    rId = prs.slides._sldIdLst[slide_index].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        # Try alternative approach
        slide_id = prs.slides._sldIdLst[slide_index]
        prs.slides._sldIdLst.remove(slide_id)
    else:
        prs.part.drop_rel(rId)
        slide_id = prs.slides._sldIdLst[slide_index]
        prs.slides._sldIdLst.remove(slide_id)
    return True


def reorder_slide(prs, from_index, to_index):
    """Move a slide from one position to another."""
    slides = prs.slides._sldIdLst
    num_slides = len(slides)

    if from_index < 0 or from_index >= num_slides:
        return False
    if to_index < 0 or to_index >= num_slides:
        return False

    el = slides[from_index]
    slides.remove(el)

    if to_index >= len(slides):
        slides.append(el)
    else:
        slides.insert(to_index, el)

    return True


def duplicate_slide(prs, slide_index):
    """Duplicate a slide."""
    from pptx.oxml.ns import qn
    from copy import deepcopy
    import lxml.etree as etree

    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    template = prs.slides[slide_index]
    slide_layout = template.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # Copy all elements from template to new slide
    for shape in template.shapes:
        el = deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)

    # Remove default placeholder shapes that came with the layout
    # (keep only the copied ones)

    return True


def update_slide_title(prs, slide_index, title_text):
    """Update the title of a specific slide."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    slide = prs.slides[slide_index]
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        return True
    return False


def update_notes(prs, slide_index, notes_text):
    """Update speaker notes for a specific slide."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    slide = prs.slides[slide_index]
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text
    return True


def add_image_to_slide(prs, slide_index, image_path, left_inches, top_inches, width_inches=None, height_inches=None):
    """Add an image to a specific slide."""
    from pptx.util import Inches

    if slide_index < 0 or slide_index >= len(prs.slides):
        return False
    if not os.path.exists(image_path):
        return False

    slide = prs.slides[slide_index]
    left = Inches(left_inches)
    top = Inches(top_inches)
    width = Inches(width_inches) if width_inches else None
    height = Inches(height_inches) if height_inches else None

    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    return True


def add_textbox_to_slide(prs, slide_index, text, left_inches, top_inches,
                          width_inches, height_inches, **kwargs):
    """Add a textbox to a specific slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    slide = prs.slides[slide_index]
    left = Inches(left_inches)
    top = Inches(top_inches)
    width = Inches(width_inches)
    height = Inches(height_inches)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    if kwargs.get('font_size'):
        run.font.size = Pt(kwargs['font_size'])
    if kwargs.get('bold'):
        run.font.bold = True
    if kwargs.get('italic'):
        run.font.italic = True
    if kwargs.get('color'):
        hex_c = kwargs['color'].lstrip('#')
        run.font.color.rgb = RGBColor(
            int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
        )
    if kwargs.get('font_name'):
        run.font.name = kwargs['font_name']

    return True


def add_table_to_slide(prs, slide_index, headers, rows,
                        left_inches=1.0, top_inches=2.5,
                        width_inches=10.0, height_inches=None):
    """Add a table to a specific slide."""
    from pptx.util import Inches, Pt

    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    slide = prs.slides[slide_index]
    num_cols = len(headers)
    num_rows = len(rows) + 1

    if not height_inches:
        height_inches = 0.5 * num_rows

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(left_inches), Inches(top_inches),
        Inches(width_inches), Inches(height_inches)
    )
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True

    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            if c < num_cols:
                table.cell(r + 1, c).text = str(val)

    return True


def add_shape_to_slide(prs, slide_index, shape_type_name, left_inches, top_inches,
                        width_inches, height_inches, **kwargs):
    """Add a shape to a specific slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    shape_types = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'oval': MSO_SHAPE.OVAL,
        'diamond': MSO_SHAPE.DIAMOND,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'right_arrow': MSO_SHAPE.RIGHT_ARROW,
        'left_arrow': MSO_SHAPE.LEFT_ARROW,
        'pentagon': MSO_SHAPE.PENTAGON,
        'hexagon': MSO_SHAPE.HEXAGON,
        'star': MSO_SHAPE.STAR_5_POINT,
        'cloud': MSO_SHAPE.CLOUD,
        'heart': MSO_SHAPE.HEART,
        'chevron': MSO_SHAPE.CHEVRON,
    }

    mso_shape = shape_types.get(shape_type_name, MSO_SHAPE.RECTANGLE)
    slide = prs.slides[slide_index]

    shape = slide.shapes.add_shape(
        mso_shape,
        Inches(left_inches), Inches(top_inches),
        Inches(width_inches), Inches(height_inches)
    )

    if kwargs.get('fill_color'):
        hex_c = kwargs['fill_color'].lstrip('#')
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
        )

    if kwargs.get('text'):
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = kwargs['text']
        if kwargs.get('font_color'):
            hex_c = kwargs['font_color'].lstrip('#')
            run.font.color.rgb = RGBColor(
                int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
            )
        if kwargs.get('font_size'):
            run.font.size = Pt(kwargs['font_size'])

    return True


def set_slide_background(prs, slide_index, color):
    """Set the background color of a slide."""
    from pptx.dml.color import RGBColor

    if slide_index < 0 or slide_index >= len(prs.slides):
        return False

    slide = prs.slides[slide_index]
    background = slide.background
    fill = background.fill
    fill.solid()
    hex_c = color.lstrip('#')
    fill.fore_color.rgb = RGBColor(
        int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
    )
    return True


def update_metadata(prs, **kwargs):
    """Update presentation core properties."""
    props = prs.core_properties
    if 'title' in kwargs:
        props.title = kwargs['title']
    if 'author' in kwargs:
        props.author = kwargs['author']
    if 'subject' in kwargs:
        props.subject = kwargs['subject']
    if 'keywords' in kwargs:
        props.keywords = kwargs['keywords']
    if 'comments' in kwargs:
        props.comments = kwargs['comments']
    return True


def add_new_slide(prs, position, layout_name, slide_spec):
    """Add a new slide at a specific position."""
    # Import create helpers
    sys.path.insert(0, os.path.dirname(__file__))
    from create_pptx import get_layout, add_content_to_textframe, add_shapes_to_slide, add_table_to_slide

    layout = get_layout(prs, layout_name)
    slide = prs.slides.add_slide(layout)

    # Set title
    if slide_spec.get('title') and slide.shapes.title:
        slide.shapes.title.text = slide_spec['title']

    # Set subtitle
    if slide_spec.get('subtitle'):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = slide_spec['subtitle']
                break

    # Content
    if slide_spec.get('content'):
        body_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_placeholder = ph
                break

        if body_placeholder and body_placeholder.has_text_frame:
            tf = body_placeholder.text_frame
            tf.clear()
            text_items = [c for c in slide_spec['content'] if c.get('type') != 'table']
            table_items = [c for c in slide_spec['content'] if c.get('type') == 'table']
            if text_items:
                add_content_to_textframe(tf, text_items)
            for table_item in table_items:
                add_table_to_slide(slide, table_item)

    # Shapes
    if slide_spec.get('shapes'):
        add_shapes_to_slide(slide, slide_spec['shapes'])

    # Notes
    if slide_spec.get('notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_spec['notes']

    # Reorder to requested position
    if position is not None and position < len(prs.slides) - 1:
        reorder_slide(prs, len(prs.slides) - 1, position)

    return True


def process_operations(input_path, output_path, operations):
    """Process a list of edit operations on a presentation."""
    ensure_dependencies()
    from pptx import Presentation

    prs = Presentation(input_path)
    results = []

    for op in operations:
        action = op.get('action', '')
        result = {"action": action, "status": "success"}

        try:
            if action == 'find_replace':
                count = find_replace(prs, op['find'], op['replace'])
                result["replacements"] = count
                if count == 0:
                    result["status"] = "no_matches"
                    result["message"] = f"No matches found for '{op['find']}'"

            elif action == 'add_slide':
                position = op.get('position')
                layout = op.get('layout', 'blank')
                success = add_new_slide(prs, position, layout, op)
                if not success:
                    result["status"] = "error"
                    result["message"] = "Failed to add slide"

            elif action == 'delete_slide':
                success = delete_slide(prs, op['slide_index'])
                if not success:
                    result["status"] = "error"
                    result["message"] = f"Invalid slide index: {op['slide_index']}"

            elif action == 'reorder_slide':
                success = reorder_slide(prs, op['from_index'], op['to_index'])
                if not success:
                    result["status"] = "error"
                    result["message"] = "Invalid slide indices"

            elif action == 'duplicate_slide':
                success = duplicate_slide(prs, op['slide_index'])
                if not success:
                    result["status"] = "error"
                    result["message"] = f"Invalid slide index: {op['slide_index']}"

            elif action == 'update_slide_title':
                success = update_slide_title(prs, op['slide_index'], op['title'])
                if not success:
                    result["status"] = "warning"
                    result["message"] = "Slide has no title placeholder"

            elif action == 'update_notes':
                success = update_notes(prs, op['slide_index'], op['notes'])
                if not success:
                    result["status"] = "error"
                    result["message"] = f"Invalid slide index: {op['slide_index']}"

            elif action == 'add_image':
                success = add_image_to_slide(
                    prs, op['slide_index'], op['image_path'],
                    op.get('left_inches', 1.0), op.get('top_inches', 1.0),
                    op.get('width_inches'), op.get('height_inches')
                )
                if not success:
                    result["status"] = "error"
                    result["message"] = "Image not found or invalid slide index"

            elif action == 'add_textbox':
                success = add_textbox_to_slide(
                    prs, op['slide_index'], op['text'],
                    op.get('left_inches', 1.0), op.get('top_inches', 1.0),
                    op.get('width_inches', 4.0), op.get('height_inches', 1.0),
                    font_size=op.get('font_size'),
                    bold=op.get('bold'),
                    italic=op.get('italic'),
                    color=op.get('color'),
                    font_name=op.get('font_name'),
                )
                if not success:
                    result["status"] = "error"
                    result["message"] = f"Invalid slide index: {op['slide_index']}"

            elif action == 'add_table':
                success = add_table_to_slide(
                    prs, op['slide_index'], op['headers'], op['rows'],
                    op.get('left_inches', 1.0), op.get('top_inches', 2.5),
                    op.get('width_inches', 10.0), op.get('height_inches')
                )
                if not success:
                    result["status"] = "error"
                    result["message"] = f"Invalid slide index: {op['slide_index']}"

            elif action == 'add_shape':
                success = add_shape_to_slide(
                    prs, op['slide_index'],
                    op.get('shape_type', 'rectangle'),
                    op.get('left_inches', 1.0), op.get('top_inches', 1.0),
                    op.get('width_inches', 4.0), op.get('height_inches', 2.0),
                    fill_color=op.get('fill_color'),
                    text=op.get('text'),
                    font_color=op.get('font_color'),
                    font_size=op.get('font_size'),
                )
                if not success:
                    result["status"] = "error"
                    result["message"] = f"Invalid slide index: {op['slide_index']}"

            elif action == 'set_background':
                success = set_slide_background(prs, op['slide_index'], op['color'])
                if not success:
                    result["status"] = "error"
                    result["message"] = f"Invalid slide index: {op['slide_index']}"

            elif action == 'update_metadata':
                update_metadata(prs, **{k: v for k, v in op.items() if k != 'action'})

            else:
                result["status"] = "error"
                result["message"] = f"Unknown action: {action}"

        except Exception as e:
            result["status"] = "error"
            result["message"] = str(e)

        results.append(result)

    os.makedirs(os.path.dirname(output_path) or '/tmp', exist_ok=True)
    prs.save(output_path)

    return {"status": "success", "output_path": output_path, "operations": results}


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: edit_pptx.py <json_operations_file_or_string>"}))
        sys.exit(1)

    input_arg = sys.argv[1]

    if os.path.isfile(input_arg):
        with open(input_arg, 'r') as f:
            spec = json.load(f)
    else:
        spec = json.loads(input_arg)

    input_path = spec.get('input_path')
    output_path = spec.get('output_path',
                           input_path.replace('.pptx', '_edited.pptx') if input_path else '/tmp/edited.pptx')
    operations = spec.get('operations', [])

    if not input_path or not os.path.exists(input_path):
        print(json.dumps({"error": f"Input file not found: {input_path}"}))
        sys.exit(1)

    result = process_operations(input_path, output_path, operations)
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
