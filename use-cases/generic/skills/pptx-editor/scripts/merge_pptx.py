#!/usr/bin/env python3
"""Merge multiple PowerPoint presentations into one."""

import json
import sys
import os
import argparse


def ensure_dependencies():
    try:
        import pptx
    except ImportError:
        os.system("pip install python-pptx")


def merge_presentations(file_paths, output_path, section_breaks=False):
    """Merge multiple .pptx files into a single presentation."""
    ensure_dependencies()
    from pptx import Presentation
    from copy import deepcopy

    if not file_paths:
        return {"error": "No files provided"}

    # Check all files exist
    for fp in file_paths:
        if not os.path.exists(fp):
            return {"error": f"File not found: {fp}"}

    # Use first presentation as base
    merged = Presentation(file_paths[0])
    base_width = merged.slide_width
    base_height = merged.slide_height

    warnings = []
    slide_counts = [len(merged.slides)]

    for file_path in file_paths[1:]:
        src = Presentation(file_path)

        # Check dimension compatibility
        if src.slide_width != base_width or src.slide_height != base_height:
            warnings.append(
                f"'{os.path.basename(file_path)}' has different dimensions "
                f"({round(src.slide_width / 914400, 2)}x{round(src.slide_height / 914400, 2)}) "
                f"vs base ({round(base_width / 914400, 2)}x{round(base_height / 914400, 2)})"
            )

        # Add section break if requested
        if section_breaks:
            # Add a blank slide with the filename as title
            blank_layout = merged.slide_layouts[6] if len(merged.slide_layouts) > 6 else merged.slide_layouts[0]
            section_slide = merged.slides.add_slide(blank_layout)
            if section_slide.shapes.title:
                section_slide.shapes.title.text = os.path.splitext(os.path.basename(file_path))[0]

        src_slide_count = 0
        for slide in src.slides:
            # Find best matching layout in merged presentation
            src_layout_name = slide.slide_layout.name
            target_layout = None
            for layout in merged.slide_layouts:
                if layout.name == src_layout_name:
                    target_layout = layout
                    break
            if not target_layout:
                target_layout = merged.slide_layouts[6] if len(merged.slide_layouts) > 6 else merged.slide_layouts[0]

            new_slide = merged.slides.add_slide(target_layout)

            # Copy all shapes from source slide
            for shape in slide.shapes:
                el = deepcopy(shape.element)
                new_slide.shapes._spTree.append(el)

            # Copy notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    new_slide.notes_slide.notes_text_frame.text = notes_text

            src_slide_count += 1

        slide_counts.append(src_slide_count)

    os.makedirs(os.path.dirname(output_path) or '/tmp', exist_ok=True)
    merged.save(output_path)

    result = {
        "status": "success",
        "output_path": output_path,
        "files_merged": len(file_paths),
        "total_slides": len(merged.slides),
        "slides_per_source": dict(zip(
            [os.path.basename(f) for f in file_paths],
            slide_counts
        )),
    }
    if warnings:
        result["warnings"] = warnings

    return result


def main():
    parser = argparse.ArgumentParser(description="Merge multiple .pptx files into one")
    parser.add_argument('--files', nargs='+', required=True, help="Input .pptx file paths")
    parser.add_argument('--output', required=True, help="Output .pptx file path")
    parser.add_argument('--section_breaks', action='store_true', default=False,
                        help="Insert section header slides between decks")

    args = parser.parse_args()

    result = merge_presentations(args.files, args.output, args.section_breaks)
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
