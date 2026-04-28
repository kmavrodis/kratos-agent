#!/usr/bin/env python3
"""Convert PowerPoint presentations to/from various formats."""

import json
import sys
import os
import subprocess
import shutil


def ensure_dependencies():
    try:
        import pptx
    except ImportError:
        os.system("pip install python-pptx")


def pptx_to_text(input_path, output_path):
    """Convert .pptx to plain text."""
    ensure_dependencies()
    from pptx import Presentation

    prs = Presentation(input_path)
    lines = []

    for i, slide in enumerate(prs.slides):
        lines.append(f"=== Slide {i + 1} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        indent = '  ' * paragraph.level
                        lines.append(f"{indent}{text}")
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    lines.append(' | '.join(row_text))
        lines.append('')

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))

    return {"status": "success", "output_path": output_path}


def pptx_to_markdown(input_path, output_path):
    """Convert .pptx to Markdown outline."""
    ensure_dependencies()
    from pptx import Presentation

    prs = Presentation(input_path)
    md_lines = []

    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {i + 1}"
        md_lines.append(f"## Slide {i + 1}: {title}")
        md_lines.append('')

        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        indent = '  ' * paragraph.level
                        md_lines.append(f"{indent}- {text}")

            if shape.has_table:
                md_lines.append('')
                for j, row in enumerate(shape.table.rows):
                    cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
                    md_lines.append('| ' + ' | '.join(cells) + ' |')
                    if j == 0:
                        md_lines.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
                md_lines.append('')

        # Notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_lines.append('')
                md_lines.append(f"> **Speaker Notes:** {notes}")

        md_lines.append('')
        md_lines.append('---')
        md_lines.append('')

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))

    return {"status": "success", "output_path": output_path}


def pptx_to_html(input_path, output_path):
    """Convert .pptx to HTML."""
    ensure_dependencies()
    from pptx import Presentation

    prs = Presentation(input_path)
    title = prs.core_properties.title or "Presentation"

    html_parts = [
        '<!DOCTYPE html>', '<html>', '<head>',
        '<meta charset="utf-8">',
        f'<title>{title}</title>',
        '<style>',
        'body { font-family: Calibri, Arial, sans-serif; max-width: 900px; margin: 0 auto; padding: 20px; }',
        '.slide { border: 1px solid #ccc; border-radius: 8px; padding: 30px; margin: 20px 0; background: #fafafa; }',
        '.slide-number { color: #888; font-size: 12px; margin-bottom: 10px; }',
        '.slide h2 { margin-top: 0; color: #2E74B5; }',
        'table { border-collapse: collapse; width: 100%; margin: 1em 0; }',
        'th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }',
        'th { background-color: #2E74B5; color: white; }',
        '.notes { background: #fff3cd; border-left: 4px solid #ffc107; padding: 10px; margin-top: 15px; font-size: 14px; }',
        'ul { margin: 5px 0; }',
        '</style>',
        '</head>', '<body>',
        f'<h1>{_escape_html(title)}</h1>',
    ]

    for i, slide in enumerate(prs.slides):
        html_parts.append(f'<div class="slide">')
        html_parts.append(f'<div class="slide-number">Slide {i + 1}</div>')

        if slide.shapes.title:
            html_parts.append(f'<h2>{_escape_html(slide.shapes.title.text)}</h2>')

        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                html_parts.append('<ul>')
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        html_parts.append(f'<li>{_escape_html(text)}</li>')
                html_parts.append('</ul>')

            if shape.has_table:
                html_parts.append('<table>')
                for j, row in enumerate(shape.table.rows):
                    html_parts.append('<tr>')
                    tag = 'th' if j == 0 else 'td'
                    for cell in row.cells:
                        html_parts.append(f'<{tag}>{_escape_html(cell.text)}</{tag}>')
                    html_parts.append('</tr>')
                html_parts.append('</table>')

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                html_parts.append(f'<div class="notes"><strong>Notes:</strong> {_escape_html(notes)}</div>')

        html_parts.append('</div>')

    html_parts.extend(['</body>', '</html>'])

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(html_parts))

    return {"status": "success", "output_path": output_path}


def pptx_to_pdf(input_path, output_path):
    """Convert .pptx to PDF using LibreOffice."""
    lo_path = shutil.which('libreoffice') or shutil.which('soffice')
    if not lo_path:
        return {
            "status": "error",
            "message": "LibreOffice not found. Install for PDF conversion: brew install --cask libreoffice (macOS) or apt install libreoffice (Linux)"
        }

    output_dir = os.path.dirname(output_path) or '/tmp'
    try:
        result = subprocess.run(
            [lo_path, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_path],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0:
            return {"status": "error", "message": result.stderr}

        expected_name = os.path.splitext(os.path.basename(input_path))[0] + '.pdf'
        expected_path = os.path.join(output_dir, expected_name)

        if expected_path != output_path and os.path.exists(expected_path):
            shutil.move(expected_path, output_path)

        return {"status": "success", "output_path": output_path}
    except subprocess.TimeoutExpired:
        return {"status": "error", "message": "PDF conversion timed out"}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def pptx_to_images(input_path, output_dir):
    """Convert .pptx to PNG images (one per slide) using LibreOffice."""
    lo_path = shutil.which('libreoffice') or shutil.which('soffice')
    if not lo_path:
        return {
            "status": "error",
            "message": "LibreOffice not found. Required for image export."
        }

    os.makedirs(output_dir, exist_ok=True)

    # First convert to PDF, then use a tool to split
    # Alternative: convert directly to images if possible
    try:
        # Convert to PDF first
        pdf_path = os.path.join(output_dir, 'temp_presentation.pdf')
        result = subprocess.run(
            [lo_path, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_path],
            capture_output=True, text=True, timeout=120
        )

        expected_name = os.path.splitext(os.path.basename(input_path))[0] + '.pdf'
        expected_path = os.path.join(output_dir, expected_name)

        if os.path.exists(expected_path):
            if expected_path != pdf_path:
                shutil.move(expected_path, pdf_path)

        # Try to convert PDF to images using pdf2image
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=200)
            image_paths = []
            for i, image in enumerate(images):
                img_path = os.path.join(output_dir, f'slide_{i + 1:03d}.png')
                image.save(img_path, 'PNG')
                image_paths.append(img_path)

            # Clean up temp PDF
            os.remove(pdf_path)

            return {
                "status": "success",
                "output_dir": output_dir,
                "images": image_paths,
                "count": len(image_paths)
            }
        except ImportError:
            return {
                "status": "partial",
                "message": "PDF created but pdf2image not available. Install with: pip install pdf2image (requires poppler)",
                "pdf_path": pdf_path
            }

    except subprocess.TimeoutExpired:
        return {"status": "error", "message": "Conversion timed out"}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def _escape_html(text):
    """Escape HTML special characters."""
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: convert_pptx.py <input_file> --to <format> [--output <output_path>]"}))
        sys.exit(1)

    input_path = sys.argv[1]
    to_format = None
    output_path = None

    if "--to" in sys.argv:
        idx = sys.argv.index("--to")
        if idx + 1 < len(sys.argv):
            to_format = sys.argv[idx + 1].lower()

    if "--output" in sys.argv:
        idx = sys.argv.index("--output")
        if idx + 1 < len(sys.argv):
            output_path = sys.argv[idx + 1]

    if not to_format:
        print(json.dumps({"error": "Specify target format with --to (pdf, txt, html, md, images)"}))
        sys.exit(1)

    if not os.path.exists(input_path):
        print(json.dumps({"error": f"File not found: {input_path}"}))
        sys.exit(1)

    base_name = os.path.splitext(os.path.basename(input_path))[0]

    converters = {
        'txt': ('text', lambda: pptx_to_text(input_path, output_path or f"/tmp/{base_name}.txt")),
        'md': ('markdown', lambda: pptx_to_markdown(input_path, output_path or f"/tmp/{base_name}.md")),
        'html': ('html', lambda: pptx_to_html(input_path, output_path or f"/tmp/{base_name}.html")),
        'pdf': ('pdf', lambda: pptx_to_pdf(input_path, output_path or f"/tmp/{base_name}.pdf")),
        'images': ('images', lambda: pptx_to_images(input_path, output_path or f"/tmp/{base_name}_slides/")),
        'png': ('images', lambda: pptx_to_images(input_path, output_path or f"/tmp/{base_name}_slides/")),
    }

    if to_format not in converters:
        print(json.dumps({"error": f"Unsupported format: {to_format}. Use: pdf, txt, html, md, images"}))
        sys.exit(1)

    _, converter_fn = converters[to_format]
    result = converter_fn()
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
