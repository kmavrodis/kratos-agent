#!/usr/bin/env python3
"""Fill placeholder values in a PowerPoint template."""

import json
import sys
import os
import re


def ensure_dependencies():
    try:
        import pptx
    except ImportError:
        os.system("pip install python-pptx")


def find_placeholders(prs):
    """Find all {{placeholder}} patterns in the presentation."""
    placeholders = set()
    pattern = re.compile(r'\{\{(\w+)\}\}')

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    matches = pattern.findall(text)
                    placeholders.update(matches)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            matches = pattern.findall(paragraph.text)
                            placeholders.update(matches)

        # Notes
        if slide.has_notes_slide:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                matches = pattern.findall(paragraph.text)
                placeholders.update(matches)

    return sorted(list(placeholders))


def replace_in_runs(paragraph, placeholder, value):
    """Replace placeholder in paragraph runs, handling split across runs."""
    target = '{{' + placeholder + '}}'
    full_text = paragraph.text

    if target not in full_text:
        return False

    # Try run-level replacement first
    for run in paragraph.runs:
        if target in run.text:
            run.text = run.text.replace(target, str(value))
            return True

    # Handle split across runs
    combined = ''.join(run.text for run in paragraph.runs)
    if target in combined:
        new_combined = combined.replace(target, str(value))
        runs = list(paragraph.runs)
        for i, run in enumerate(runs):
            if i == 0:
                run.text = new_combined
            else:
                run.text = ''
        return True

    return False


def fill_template(template_path, values, output_path):
    """Fill all placeholders in a template presentation."""
    ensure_dependencies()
    from pptx import Presentation

    prs = Presentation(template_path)

    # Find all placeholders
    all_placeholders = find_placeholders(prs)

    filled = []
    missing_values = []
    unused_values = []

    for ph in all_placeholders:
        if ph in values:
            filled.append(ph)
        else:
            missing_values.append(ph)

    for key in values:
        if key not in all_placeholders:
            unused_values.append(key)

    # Do replacements
    for placeholder, value in values.items():
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        replace_in_runs(paragraph, placeholder, value)

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                replace_in_runs(paragraph, placeholder, value)

            # Notes
            if slide.has_notes_slide:
                for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                    replace_in_runs(paragraph, placeholder, value)

    os.makedirs(os.path.dirname(output_path) or '/tmp', exist_ok=True)
    prs.save(output_path)

    result = {
        "status": "success",
        "output_path": output_path,
        "placeholders_found": all_placeholders,
        "filled": filled,
    }
    if missing_values:
        result["missing_values"] = missing_values
        result["warning"] = f"No values provided for: {', '.join(missing_values)}"
    if unused_values:
        result["unused_values"] = unused_values

    return result


def main():
    if len(sys.argv) < 2:
        print(json.dumps({
            "error": "Usage: template_fill.py <template.pptx> --values '{\"key\": \"value\"}' --output <output.pptx>"
        }))
        sys.exit(1)

    ensure_dependencies()

    template_path = sys.argv[1]

    if not os.path.exists(template_path):
        print(json.dumps({"error": f"Template not found: {template_path}"}))
        sys.exit(1)

    # Scan-only mode
    if "--scan" in sys.argv:
        from pptx import Presentation
        prs = Presentation(template_path)
        placeholders = find_placeholders(prs)
        print(json.dumps({"placeholders": placeholders, "count": len(placeholders)}, indent=2))
        return

    values = {}
    output_path = template_path.replace('.pptx', '_filled.pptx')

    if "--values" in sys.argv:
        idx = sys.argv.index("--values")
        if idx + 1 < len(sys.argv):
            values_arg = sys.argv[idx + 1]
            if os.path.isfile(values_arg):
                with open(values_arg, 'r') as f:
                    values = json.load(f)
            else:
                values = json.loads(values_arg)

    if "--output" in sys.argv:
        idx = sys.argv.index("--output")
        if idx + 1 < len(sys.argv):
            output_path = sys.argv[idx + 1]

    result = fill_template(template_path, values, output_path)
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
