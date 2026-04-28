#!/usr/bin/env python3
"""Read and extract content from PowerPoint presentations (.pptx)."""

import json
import sys
import os


def ensure_dependencies():
    try:
        import pptx
    except ImportError:
        os.system("pip install python-pptx")


def extract_text(prs_path):
    """Extract plain text from all slides."""
    from pptx import Presentation
    prs = Presentation(prs_path)

    slides_text = []
    for i, slide in enumerate(prs.slides):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    lines.append(' | '.join(row_text))
        slides_text.append({
            "slide": i + 1,
            "text": '\n'.join(lines)
        })

    return {"slides": slides_text}


def extract_outline(prs_path):
    """Extract titles and bullet points as an outline."""
    from pptx import Presentation
    prs = Presentation(prs_path)

    outline = []
    for i, slide in enumerate(prs.slides):
        slide_info = {"slide": i + 1, "title": "", "bullets": []}

        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text

        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_info["bullets"].append({
                            "text": text,
                            "level": paragraph.level,
                        })

        outline.append(slide_info)

    return {"outline": outline, "slide_count": len(outline)}


def extract_notes(prs_path):
    """Extract speaker notes from all slides."""
    from pptx import Presentation
    prs = Presentation(prs_path)

    notes = []
    for i, slide in enumerate(prs.slides):
        note_text = ""
        if slide.has_notes_slide:
            note_text = slide.notes_slide.notes_text_frame.text.strip()

        title = slide.shapes.title.text if slide.shapes.title else f"Slide {i + 1}"
        notes.append({
            "slide": i + 1,
            "title": title,
            "notes": note_text,
        })

    return {"notes": notes}


def extract_structured(prs_path):
    """Extract full structured content with formatting details."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation(prs_path)

    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide": i + 1,
            "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "shapes": [],
        }

        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text

        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left_inches": round(shape.left / 914400, 2) if shape.left else None,
                "top_inches": round(shape.top / 914400, 2) if shape.top else None,
                "width_inches": round(shape.width / 914400, 2) if shape.width else None,
                "height_inches": round(shape.height / 914400, 2) if shape.height else None,
            }

            if shape.has_text_frame:
                paragraphs = []
                for p in shape.text_frame.paragraphs:
                    runs = []
                    for run in p.runs:
                        run_info = {
                            "text": run.text,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                        }
                        if run.font.size:
                            run_info["font_size"] = run.font.size.pt
                        if run.font.name:
                            run_info["font_name"] = run.font.name
                        if run.font.color and run.font.color.rgb:
                            run_info["color"] = str(run.font.color.rgb)
                        runs.append(run_info)
                    paragraphs.append({
                        "text": p.text,
                        "level": p.level,
                        "alignment": str(p.alignment) if p.alignment else None,
                        "runs": runs,
                    })
                shape_info["paragraphs"] = paragraphs

            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text for cell in row.cells])
                shape_info["table"] = table_data

            slide_info["shapes"].append(shape_info)

        # Notes
        if slide.has_notes_slide:
            slide_info["notes"] = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append(slide_info)

    return {
        "slide_count": len(slides_data),
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "slides": slides_data,
    }


def extract_tables(prs_path):
    """Extract all tables from the presentation."""
    from pptx import Presentation
    prs = Presentation(prs_path)

    tables = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_table:
                table_data = {
                    "slide": i + 1,
                    "rows": []
                }
                for row in shape.table.rows:
                    table_data["rows"].append([cell.text.strip() for cell in row.cells])
                tables.append(table_data)

    return {"tables": tables, "count": len(tables)}


def extract_images(prs_path, output_dir="/tmp/pptx_images"):
    """Extract embedded images to output directory."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(prs_path)

    image_paths = []
    image_count = 0

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                image = shape.image
                ext = image.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'
                filename = f"slide{i + 1}_image{image_count}.{ext}"
                filepath = os.path.join(output_dir, filename)
                with open(filepath, 'wb') as f:
                    f.write(image.blob)
                image_paths.append(filepath)

    return {"images": image_paths, "count": image_count}


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: read_pptx.py <file.pptx> [--mode text|structured|notes|outline|images|tables]"}))
        sys.exit(1)

    ensure_dependencies()

    prs_path = sys.argv[1]
    mode = "text"

    if "--mode" in sys.argv:
        idx = sys.argv.index("--mode")
        if idx + 1 < len(sys.argv):
            mode = sys.argv[idx + 1]

    if not os.path.exists(prs_path):
        print(json.dumps({"error": f"File not found: {prs_path}"}))
        sys.exit(1)

    try:
        if mode == "text":
            result = extract_text(prs_path)
        elif mode == "structured":
            result = extract_structured(prs_path)
        elif mode == "notes":
            result = extract_notes(prs_path)
        elif mode == "outline":
            result = extract_outline(prs_path)
        elif mode == "images":
            output_dir = "/tmp/pptx_images"
            if "--output" in sys.argv:
                idx = sys.argv.index("--output")
                if idx + 1 < len(sys.argv):
                    output_dir = sys.argv[idx + 1]
            result = extract_images(prs_path, output_dir)
        elif mode == "tables":
            result = extract_tables(prs_path)
        else:
            result = {"error": f"Unknown mode: {mode}. Use: text, structured, notes, outline, images, tables"}

        print(json.dumps(result, indent=2, default=str))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)


if __name__ == '__main__':
    main()
