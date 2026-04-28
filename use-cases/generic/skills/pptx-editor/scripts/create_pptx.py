#!/usr/bin/env python3
"""Create a PowerPoint presentation from a JSON content specification."""

import json
import sys
import os
import re


def ensure_dependencies():
    try:
        import pptx
    except ImportError:
        os.system("pip install python-pptx")


def apply_text_formatting(run, item):
    """Apply font settings to a text run."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    if item.get('bold'):
        run.font.bold = True
    if item.get('italic'):
        run.font.italic = True
    if item.get('underline'):
        run.font.underline = True
    if item.get('font_name'):
        run.font.name = item['font_name']
    if item.get('font_size'):
        run.font.size = Pt(item['font_size'])
    if item.get('color'):
        hex_c = item['color'].lstrip('#')
        run.font.color.rgb = RGBColor(
            int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
        )


def set_paragraph_alignment(paragraph, alignment_str):
    """Set paragraph alignment."""
    from pptx.enum.text import PP_ALIGN
    alignments = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY,
    }
    if alignment_str and alignment_str.lower() in alignments:
        paragraph.alignment = alignments[alignment_str.lower()]


def get_layout(prs, layout_name):
    """Get a slide layout by name or index, with fallback."""
    layout_map = {
        'title': 0,
        'title_content': 1,
        'section_header': 2,
        'two_column': 3,
        'comparison': 4,
        'title_only': 5,
        'blank': 6,
        'content_caption': 7,
    }

    # Try by our alias first
    idx = layout_map.get(layout_name)
    if idx is not None and idx < len(prs.slide_layouts):
        return prs.slide_layouts[idx]

    # Try by actual layout name
    for layout in prs.slide_layouts:
        if layout.name.lower().replace(' ', '_') == layout_name.lower().replace(' ', '_'):
            return layout

    # Fallback to blank or first available
    blank_idx = layout_map.get('blank', 6)
    if blank_idx < len(prs.slide_layouts):
        return prs.slide_layouts[blank_idx]
    return prs.slide_layouts[0]


def add_content_to_textframe(tf, content_items):
    """Add content items (text, bullets, numbered lists) to a text frame."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    first = True
    for item in content_items:
        item_type = item.get('type', 'text')

        if item_type == 'text':
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = item.get('text', '')
            apply_text_formatting(run, item)
            if item.get('alignment'):
                set_paragraph_alignment(p, item['alignment'])

        elif item_type in ('bullet_list', 'numbered_list'):
            items = item.get('items', [])
            level = item.get('level', 0)
            for bullet_text in items:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.level = level
                run = p.add_run()
                run.text = str(bullet_text)
                if item.get('font_size'):
                    run.font.size = Pt(item['font_size'])
                if item.get('color'):
                    hex_c = item['color'].lstrip('#')
                    run.font.color.rgb = RGBColor(
                        int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                    )

        elif item_type == 'table':
            # Tables are handled separately as shapes, skip here
            pass


def add_table_to_slide(slide, table_item):
    """Add a table shape to a slide."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    headers = table_item.get('headers', [])
    rows_data = table_item.get('rows', [])
    num_cols = len(headers)
    num_rows = len(rows_data) + (1 if headers else 0)

    left = Inches(table_item.get('left_inches', 1.0))
    top = Inches(table_item.get('top_inches', 2.0))
    width = Inches(table_item.get('width_inches', 11.0))
    height = Inches(table_item.get('height_inches', 0.5 * num_rows))

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    font_size = table_item.get('font_size', 12)

    # Headers
    if headers:
        for i, header_text in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header_text)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(font_size)
                    if table_item.get('header_font_color'):
                        hex_c = table_item['header_font_color'].lstrip('#')
                        run.font.color.rgb = RGBColor(
                            int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                        )

            if table_item.get('header_bg_color'):
                hex_c = table_item['header_bg_color'].lstrip('#')
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = RGBColor(
                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                )

    # Data rows
    start_row = 1 if headers else 0
    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            if c < num_cols:
                cell = table.cell(start_row + r, c)
                cell.text = str(cell_text)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)

    return table_shape


def add_shapes_to_slide(slide, shapes):
    """Add custom shapes to a slide."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    for shape_def in shapes:
        shape_type = shape_def.get('type', '')

        if shape_type == 'textbox':
            left = Inches(shape_def.get('left_inches', 1.0))
            top = Inches(shape_def.get('top_inches', 1.0))
            width = Inches(shape_def.get('width_inches', 5.0))
            height = Inches(shape_def.get('height_inches', 1.0))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = shape_def.get('text', '')
            apply_text_formatting(run, shape_def)
            if shape_def.get('alignment'):
                set_paragraph_alignment(p, shape_def['alignment'])

        elif shape_type == 'image':
            img_path = shape_def.get('path', '')
            if os.path.exists(img_path):
                left = Inches(shape_def.get('left_inches', 1.0))
                top = Inches(shape_def.get('top_inches', 1.0))
                width = Inches(shape_def.get('width_inches', 6.0)) if shape_def.get('width_inches') else None
                height = Inches(shape_def.get('height_inches')) if shape_def.get('height_inches') else None
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)

        elif shape_type == 'rectangle':
            left = Inches(shape_def.get('left_inches', 1.0))
            top = Inches(shape_def.get('top_inches', 1.0))
            width = Inches(shape_def.get('width_inches', 4.0))
            height = Inches(shape_def.get('height_inches', 2.0))

            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

            if shape_def.get('fill_color'):
                hex_c = shape_def['fill_color'].lstrip('#')
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                )
            else:
                shape.fill.background()

            if shape_def.get('border_color'):
                hex_c = shape_def['border_color'].lstrip('#')
                shape.line.color.rgb = RGBColor(
                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                )
            if shape_def.get('border_width'):
                shape.line.width = Pt(shape_def['border_width'])

            if shape_def.get('text'):
                shape.text_frame.word_wrap = True
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = shape_def['text']
                if shape_def.get('font_color'):
                    hex_c = shape_def['font_color'].lstrip('#')
                    run.font.color.rgb = RGBColor(
                        int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                    )
                if shape_def.get('font_size'):
                    run.font.size = Pt(shape_def['font_size'])

        elif shape_type == 'oval':
            left = Inches(shape_def.get('left_inches', 1.0))
            top = Inches(shape_def.get('top_inches', 1.0))
            width = Inches(shape_def.get('width_inches', 3.0))
            height = Inches(shape_def.get('height_inches', 3.0))

            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)

            if shape_def.get('fill_color'):
                hex_c = shape_def['fill_color'].lstrip('#')
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                )

        elif shape_type == 'line':
            from pptx.util import Inches
            start_x = Inches(shape_def.get('start_x', 1.0))
            start_y = Inches(shape_def.get('start_y', 1.0))
            end_x = Inches(shape_def.get('end_x', 5.0))
            end_y = Inches(shape_def.get('end_y', 1.0))
            width_val = end_x - start_x
            height_val = end_y - start_y

            connector = slide.shapes.add_connector(
                1, start_x, start_y, end_x, end_y  # MSO_CONNECTOR.STRAIGHT = 1
            )
            if shape_def.get('color'):
                hex_c = shape_def['color'].lstrip('#')
                connector.line.color.rgb = RGBColor(
                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                )
            if shape_def.get('width'):
                connector.line.width = Pt(shape_def['width'])

        elif shape_type == 'table':
            add_table_to_slide(slide, shape_def)


def create_presentation(spec):
    """Create a .pptx presentation from a content specification dict."""
    ensure_dependencies()
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # Slide dimensions
    if spec.get('slide_width_inches'):
        prs.slide_width = Inches(spec['slide_width_inches'])
    if spec.get('slide_height_inches'):
        prs.slide_height = Inches(spec['slide_height_inches'])

    # Metadata
    meta = spec.get('metadata', {})
    if meta.get('title'):
        prs.core_properties.title = meta['title']
    if meta.get('author'):
        prs.core_properties.author = meta['author']
    if meta.get('subject'):
        prs.core_properties.subject = meta['subject']

    # Slides
    for slide_spec in spec.get('slides', []):
        layout_name = slide_spec.get('layout', 'blank')
        layout = get_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        # Title
        if slide_spec.get('title') and slide.shapes.title:
            slide.shapes.title.text = slide_spec['title']

        # Subtitle (for title and section_header layouts)
        if slide_spec.get('subtitle'):
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:  # subtitle placeholder
                    ph.text = slide_spec['subtitle']
                    break

        # Content area (for title_content layout)
        if slide_spec.get('content'):
            body_placeholder = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body_placeholder = ph
                    break

            if body_placeholder and body_placeholder.has_text_frame:
                tf = body_placeholder.text_frame
                tf.clear()

                # Separate table items from text items
                text_items = [c for c in slide_spec['content'] if c.get('type') != 'table']
                table_items = [c for c in slide_spec['content'] if c.get('type') == 'table']

                if text_items:
                    add_content_to_textframe(tf, text_items)

                for table_item in table_items:
                    add_table_to_slide(slide, table_item)

        # Two-column layout
        if slide_spec.get('left_content') or slide_spec.get('right_content'):
            placeholders = list(slide.placeholders)
            # Usually idx 1 = left, idx 2 = right
            left_ph = None
            right_ph = None
            for ph in placeholders:
                if ph.placeholder_format.idx == 1:
                    left_ph = ph
                elif ph.placeholder_format.idx == 2:
                    right_ph = ph

            if left_ph and slide_spec.get('left_content') and left_ph.has_text_frame:
                left_ph.text_frame.clear()
                add_content_to_textframe(left_ph.text_frame, slide_spec['left_content'])
            if right_ph and slide_spec.get('right_content') and right_ph.has_text_frame:
                right_ph.text_frame.clear()
                add_content_to_textframe(right_ph.text_frame, slide_spec['right_content'])

        # Custom shapes
        if slide_spec.get('shapes'):
            add_shapes_to_slide(slide, slide_spec['shapes'])

        # Speaker notes
        if slide_spec.get('notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_spec['notes']

    # Save
    output_path = spec.get('output_path', '/tmp/presentation.pptx')
    os.makedirs(os.path.dirname(output_path) or '/tmp', exist_ok=True)
    prs.save(output_path)

    return {
        "status": "success",
        "output_path": output_path,
        "slide_count": len(prs.slides),
    }


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: create_pptx.py <json_spec_file_or_json_string>"}))
        sys.exit(1)

    input_arg = sys.argv[1]

    if os.path.isfile(input_arg):
        with open(input_arg, 'r') as f:
            spec = json.load(f)
    else:
        spec = json.loads(input_arg)

    result = create_presentation(spec)
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
